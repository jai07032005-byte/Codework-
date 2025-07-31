# pages/3_📄_Document_Q&A.py
import streamlit as st
import google.generativeai as genai
import os
from dotenv import load_dotenv

# Import the necessary libraries to read different file types
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import json # To format chat history for the prompt

# --- Load the Gemini API Key ---
try:
    load_dotenv()
    # Ensure the API key is loaded and configured
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        GEMINI_CONFIGURED = False
    else:
        genai.configure(api_key=api_key)
        GEMINI_CONFIGURED = True
except Exception as e:
    GEMINI_CONFIGURED = False
    print(f"Error configuring Gemini: {e}")

# --- Helper Function for CSS ---
def load_css():
    """Injects custom CSS for an animated gradient background and other UI enhancements."""
    st.markdown("""
        <style>
            @keyframes gradient_animation {
                0% { background-position: 0% 50%; }
                50% { background-position: 100% 50%; }
                100% { background-position: 0% 50%; }
            }
            .stApp {
                background: linear-gradient(-45deg, #1d2b64, #2c3e50, #4a5a7e, #f8cdda);
                background-size: 400% 400%;
                animation: gradient_animation 18s ease infinite;
                color: white;
            }
            .stButton>button {
                border-radius: 20px;
                border: 1px solid #f8cdda;
            }
            .stExpander {
                border-color: #f8cdda !important;
            }
        </style>
    """, unsafe_allow_html=True)

# --- Core Text Extraction Function (Remains the same, but with better error handling) ---
def extract_text_from_file(uploaded_file):
    """Extracts text from PDF, DOCX, PPTX, and TXT files."""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    try:
        if file_extension == ".pdf":
            pdf_reader = PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
        elif file_extension == ".docx":
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_extension == ".pptx":
            pres = Presentation(uploaded_file)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_extension == ".txt":
            text = uploaded_file.read().decode("utf-8")
        
        if not text.strip():
            return "Error: No text could be extracted from the document. It might be empty or scanned."
        return text
    except Exception as e:
        return f"Error reading file: {e}"

# --- Gemini API Functions (Updated for Chat) ---
def summarize_text_with_gemini(text):
    if not GEMINI_CONFIGURED: return "Error: Gemini API not configured."
    model = genai.GenerativeModel('gemini-1.5-flash')
    prompt = f"You are an expert summarizer. Please provide a detailed yet concise summary of the following document text, highlighting the key points, main arguments, and any conclusions.\n\nDOCUMENT TEXT:\n---\n{text}\n---"
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"An error occurred with the Gemini API: {e}"

def get_answer_from_gemini(document_text, chat_history, user_question):
    if not GEMINI_CONFIGURED: return "Error: Gemini API not configured."
    model = genai.GenerativeModel('gemini-1.5-flash')
    
    # Format the chat history for the prompt
    history_str = json.dumps(chat_history, indent=2)

    prompt = f"""
    You are a helpful AI assistant specialized in analyzing documents.
    Your task is to answer the user's question based *only* on the provided document text.
    Consider the previous conversation for context. If the answer is not in the document, state that clearly.

    **DOCUMENT TEXT:**
    ---
    {document_text}
    ---

    **CONVERSATION HISTORY:**
    ---
    {history_str}
    ---
    
    **USER'S LATEST QUESTION:**
    {user_question}
    """
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"An error occurred with the Gemini API: {e}"

# --- Streamlit Page Setup ---
st.set_page_config(page_title="Document Q&A", page_icon="📄", layout="wide")
load_css()

st.title("📄 Document Analysis Suite")
st.write("Upload a document to summarize it or chat about its contents.")

# --- Initialize Session State ---
if "document_text" not in st.session_state:
    st.session_state.document_text = None
    st.session_state.document_name = None
    st.session_state.chat_messages = []

# --- Main App Logic ---
if not GEMINI_CONFIGURED:
    st.error("CRITICAL ERROR: The Gemini API key is not configured. Please add it to your .env file.")
else:
    # Create two tabs for the different tools
    tab1, tab2 = st.tabs(["📄 Document Summarizer", "💬 Document Q&A"])

    with tab1:
        st.header("Summarize Any Document")
        summary_file = st.file_uploader(
            "Upload a document to get its summary",
            type=["pdf", "docx", "pptx", "txt"],
            key="summary_uploader"
        )
        if summary_file:
            if st.button("✨ Generate Summary", type="primary"):
                with st.spinner("Reading and summarizing your document..."):
                    extracted_text = extract_text_from_file(summary_file)
                    if "Error:" not in extracted_text:
                        summary = summarize_text_with_gemini(extracted_text)
                        st.subheader(f"Summary of `{summary_file.name}`")
                        st.markdown(summary)
                    else:
                        st.error(extracted_text)

    with tab2:
        st.header("Chat with Your Document")
        
        # Area for file upload and clearing session
        col1, col2 = st.columns([3, 1])
        with col1:
             qa_file = st.file_uploader(
                "Upload a document to start chatting",
                type=["pdf", "docx", "pptx", "txt"],
                key="qa_uploader"
            )
        with col2:
            st.write("") # for vertical alignment
            st.write("")
            if st.button("🗑️ Clear & Start Over"):
                st.session_state.document_text = None
                st.session_state.document_name = None
                st.session_state.chat_messages = []
                st.rerun()

        # Logic to process the uploaded file for the chat
        if qa_file and (st.session_state.document_name != qa_file.name):
            with st.spinner(f"Reading `{qa_file.name}`..."):
                st.session_state.document_text = extract_text_from_file(qa_file)
                st.session_state.document_name = qa_file.name
                st.session_state.chat_messages = [] # Reset chat history for new doc

        # Main chat interface
        if st.session_state.document_text:
            if "Error:" in st.session_state.document_text:
                st.error(f"Failed to process `{st.session_state.document_name}`. {st.session_state.document_text}")
            else:
                st.success(f"Ready to chat about **{st.session_state.document_name}**!")
                
                with st.expander("View Extracted Text"):
                    st.text(st.session_state.document_text[:2000] + "...") # Show a preview

                # Display existing chat messages
                for message in st.session_state.chat_messages:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"])

                # Get new user input
                if prompt := st.chat_input("Ask a question about your document..."):
                    # Add user message to chat history
                    st.session_state.chat_messages.append({"role": "user", "content": prompt})
                    with st.chat_message("user"):
                        st.markdown(prompt)

                    # Get and display model response
                    with st.chat_message("assistant"):
                        with st.spinner("Analyzing..."):
                            response = get_answer_from_gemini(
                                st.session_state.document_text,
                                st.session_state.chat_messages,
                                prompt
                            )
                            st.markdown(response)
                    # Add model response to chat history
                    st.session_state.chat_messages.append({"role": "assistant", "content": response})
        else:
            st.info("Please upload a document to begin the Q&A session.")